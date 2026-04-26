import io
import math
import re
from collections import Counter


STOPWORDS = {
    "a",
    "about",
    "above",
    "after",
    "again",
    "against",
    "all",
    "also",
    "am",
    "an",
    "and",
    "any",
    "are",
    "as",
    "at",
    "be",
    "because",
    "been",
    "before",
    "being",
    "below",
    "between",
    "both",
    "but",
    "by",
    "can",
    "could",
    "did",
    "do",
    "does",
    "doing",
    "down",
    "during",
    "each",
    "few",
    "for",
    "from",
    "further",
    "had",
    "has",
    "have",
    "having",
    "he",
    "her",
    "here",
    "hers",
    "herself",
    "him",
    "himself",
    "his",
    "how",
    "i",
    "if",
    "in",
    "into",
    "is",
    "it",
    "its",
    "itself",
    "just",
    "me",
    "more",
    "most",
    "my",
    "myself",
    "no",
    "nor",
    "not",
    "now",
    "of",
    "off",
    "on",
    "once",
    "only",
    "or",
    "other",
    "our",
    "ours",
    "ourselves",
    "out",
    "over",
    "own",
    "same",
    "she",
    "should",
    "so",
    "some",
    "such",
    "than",
    "that",
    "the",
    "their",
    "theirs",
    "them",
    "themselves",
    "then",
    "there",
    "these",
    "they",
    "this",
    "those",
    "through",
    "to",
    "too",
    "under",
    "until",
    "up",
    "very",
    "was",
    "we",
    "were",
    "what",
    "when",
    "where",
    "which",
    "while",
    "who",
    "whom",
    "why",
    "with",
    "would",
    "you",
    "your",
    "yours",
    "yourself",
    "yourselves",
}

GENERIC_TITLES = {
    "title slide",
    "overview",
    "summary",
    "contents",
    "introduction",
    "agenda",
    "outline",
    "learning objectives",
}

SUMMARY_FILLER_PREFIXES = (
    "this slide explains",
    "this slide describes",
    "this slide covers",
    "the slide explains",
    "the slide describes",
    "the slide covers",
    "this topic explains",
    "this topic describes",
    "this topic covers",
)

SENTENCE_SPLIT_PATTERN = re.compile(r"(?<=[.!?])\s+|\n+")
WORD_PATTERN = re.compile(r"\b[a-zA-Z][a-zA-Z0-9'-]{2,}\b")
BULLET_PREFIX_PATTERN = re.compile(r"^[\-\*\d\.\)\(]+\s*")
PARENS_PATTERN = re.compile(r"\([^)]*\)")


def normalize_text(text: str) -> str:
    safe_text = "".join(
        character
        for character in str(text or "")
        if character in "\n\t" or character.isprintable()
    )
    cleaned_lines = [re.sub(r"\s+", " ", line).strip() for line in safe_text.splitlines()]
    return "\n".join(line for line in cleaned_lines if line).strip()


def clean_fragment(text: str) -> str:
    normalized = normalize_text(text)
    normalized = BULLET_PREFIX_PATTERN.sub("", normalized)
    normalized = normalized.strip(" -:;,.")
    return re.sub(r"\s+", " ", normalized).strip()


def split_sentences(text: str) -> list[str]:
    normalized = normalize_text(text)
    return [
        clean_fragment(sentence)
        for sentence in SENTENCE_SPLIT_PATTERN.split(normalized)
        if clean_fragment(sentence)
    ]


def tokenize(text: str) -> list[str]:
    return [match.group(0).lower() for match in WORD_PATTERN.finditer(text)]


def filtered_tokens(text: str, excluded: set[str] | None = None) -> list[str]:
    excluded = excluded or set()
    return [
        token
        for token in tokenize(text)
        if token not in STOPWORDS and token not in excluded
    ]


def canonical_sentence_key(text: str) -> str:
    tokens = filtered_tokens(text)
    return " ".join(tokens[:18])


def is_generic_title(title: str) -> bool:
    normalized = clean_fragment(title).lower()
    return (
        not normalized
        or normalized in GENERIC_TITLES
        or normalized.startswith("slide ")
        or normalized.startswith("lecture ")
        or normalized.startswith("unit ")
        or normalized.startswith("chapter ")
    )


def dedupe_preserve_order(items: list[str]) -> list[str]:
    seen: set[str] = set()
    unique_items: list[str] = []

    for item in items:
        normalized = clean_fragment(item)
        if not normalized:
            continue

        marker = normalized.lower()
        if marker in seen:
            continue

        seen.add(marker)
        unique_items.append(normalized)

    return unique_items


def build_keyword_counter(text: str, excluded: set[str] | None = None) -> Counter[str]:
    return Counter(filtered_tokens(text, excluded=excluded))


def get_ranked_keywords(text: str, topic: str | None = None, limit: int = 8) -> list[str]:
    topic_tokens = set(tokenize(topic or ""))
    frequencies = build_keyword_counter(text, excluded=topic_tokens)
    ranked: list[str] = []

    for word, _count in frequencies.most_common(limit * 3):
        if len(word) < 4:
            continue
        ranked.append(word)
        if len(ranked) >= limit:
            break

    return ranked


def cosine_similarity(left: Counter[str], right: Counter[str]) -> float:
    if not left or not right:
        return 0.0

    numerator = sum(left[token] * right.get(token, 0) for token in left)
    left_norm = math.sqrt(sum(value * value for value in left.values()))
    right_norm = math.sqrt(sum(value * value for value in right.values()))

    if left_norm == 0 or right_norm == 0:
        return 0.0

    return numerator / (left_norm * right_norm)


def jaccard_similarity(left: set[str], right: set[str]) -> float:
    if not left or not right:
        return 0.0

    union = left | right
    if not union:
        return 0.0

    return len(left & right) / len(union)


def sentence_similarity(left: str, right: str) -> float:
    left_tokens = set(filtered_tokens(left))
    right_tokens = set(filtered_tokens(right))
    if not left_tokens or not right_tokens:
        return 0.0

    overlap = left_tokens & right_tokens
    jaccard = jaccard_similarity(left_tokens, right_tokens)
    containment = len(overlap) / max(1, min(len(left_tokens), len(right_tokens)))
    return max(jaccard, containment * 0.92)


def sentence_information_score(sentence: str) -> float:
    tokens = filtered_tokens(sentence)
    if not tokens:
        return 0.0

    long_word_bonus = sum(1 for word in tokens if len(word) >= 8) * 0.2
    return len(tokens) + long_word_bonus


def merge_similar_sentences(sentences: list[str], threshold: float = 0.74) -> list[str]:
    unique_sentences: list[str] = []
    canonical_indexes: dict[str, int] = {}

    for sentence in sentences:
        normalized = clean_fragment(sentence)
        if not normalized:
            continue

        canonical_key = canonical_sentence_key(normalized) or normalized.lower()
        if canonical_key in canonical_indexes:
            existing_index = canonical_indexes[canonical_key]
            if sentence_information_score(normalized) > sentence_information_score(
                unique_sentences[existing_index]
            ):
                unique_sentences[existing_index] = normalized
            continue

        matched_index = next(
            (
                index
                for index, existing in enumerate(unique_sentences)
                if sentence_similarity(normalized, existing) >= threshold
            ),
            None,
        )

        if matched_index is None:
            canonical_indexes[canonical_key] = len(unique_sentences)
            unique_sentences.append(normalized)
            continue

        if sentence_information_score(normalized) > sentence_information_score(
            unique_sentences[matched_index]
        ):
            unique_sentences[matched_index] = normalized

    return unique_sentences


def shorten_phrase(text: str, max_words: int = 18) -> str:
    words = clean_fragment(text).split()
    if len(words) <= max_words:
        return " ".join(words)
    return f"{' '.join(words[:max_words])}..."


def ensure_sentence(text: str) -> str:
    cleaned = clean_fragment(text)
    if not cleaned:
        return ""
    if cleaned.endswith((".", "!", "?")):
        return cleaned
    return f"{cleaned}."


def compress_sentence(sentence: str, topic: str | None = None, max_words: int = 16) -> str:
    cleaned = clean_fragment(sentence)
    if not cleaned:
        return ""

    cleaned = PARENS_PATTERN.sub("", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" -:;,.")

    topic_text = clean_fragment(topic or "")
    if topic_text and cleaned.lower().startswith(topic_text.lower()):
        cleaned = cleaned[len(topic_text):].lstrip(" -:;,")

    lowered = cleaned.lower()
    for prefix in SUMMARY_FILLER_PREFIXES:
        if lowered.startswith(prefix):
            cleaned = cleaned[len(prefix):].lstrip(" -:;,")
            lowered = cleaned.lower()
            break

    if not cleaned:
        cleaned = clean_fragment(sentence)

    return shorten_phrase(cleaned, max_words=max_words)


def score_sentence(
    sentence: str,
    keyword_weights: Counter[str],
    section_tokens: set[str] | None = None,
    tfidf_weights: Counter[str] | None = None,
) -> float:
    words = filtered_tokens(sentence)
    if not words:
        return 0.0

    section_tokens = section_tokens or set()
    tfidf_weights = tfidf_weights or Counter()
    score = 0.0

    for word in words:
        score += keyword_weights.get(word, 0) + 1
        score += tfidf_weights.get(word, 0) * 0.4
        if word in section_tokens:
            score += 0.9

    length_penalty = 0.82 if len(words) < 5 else 1.0
    return (score / len(words)) * length_penalty


def build_ranked_summary(
    sentences: list[str],
    topic: str | None = None,
    max_sentences: int = 2,
    tfidf_weights: Counter[str] | None = None,
) -> str:
    cleaned_sentences = merge_similar_sentences(sentences)
    if not cleaned_sentences:
        return ""

    section_tokens = set(filtered_tokens(topic or ""))
    keyword_weights = build_keyword_counter(" ".join(cleaned_sentences), excluded=section_tokens)
    scored = [
        (
            index,
            score_sentence(
                sentence,
                keyword_weights,
                section_tokens=section_tokens,
                tfidf_weights=tfidf_weights,
            ),
        )
        for index, sentence in enumerate(cleaned_sentences)
    ]
    selected_indexes = sorted(
        index
        for index, _score in sorted(scored, key=lambda item: item[1], reverse=True)[:max_sentences]
    )
    selected = [cleaned_sentences[index] for index in selected_indexes]
    return " ".join(selected)


def build_key_points(
    candidates: list[str],
    topic: str | None = None,
    limit: int = 4,
    tfidf_weights: Counter[str] | None = None,
) -> list[str]:
    if not candidates:
        return []

    cleaned_candidates = merge_similar_sentences(candidates, threshold=0.72)
    title_tokens = set(filtered_tokens(topic or ""))
    keyword_weights = build_keyword_counter(" ".join(cleaned_candidates), excluded=title_tokens)
    ranked = sorted(
        enumerate(cleaned_candidates),
        key=lambda item: score_sentence(
            item[1],
            keyword_weights,
            section_tokens=title_tokens,
            tfidf_weights=tfidf_weights,
        ),
        reverse=True,
    )

    key_points: list[str] = []
    for index, candidate in ranked:
        compressed = compress_sentence(candidate, topic=topic, max_words=16)
        if not compressed:
            continue
        if any(sentence_similarity(compressed, existing) >= 0.8 for existing in key_points):
            continue
        key_points.append(compressed)
        if len(key_points) >= limit:
            break

    if not key_points:
        for candidate in cleaned_candidates[:limit]:
            compressed = compress_sentence(candidate, topic=topic, max_words=16)
            if compressed:
                key_points.append(compressed)

    ordered_points = [
        key_points_item
        for _, key_points_item in sorted(
            (
                (
                    next(
                        (
                            original_index
                            for original_index, original in enumerate(cleaned_candidates)
                            if sentence_similarity(original, point) >= 0.65
                        ),
                        index,
                    ),
                    point,
                )
                for index, point in enumerate(key_points)
            ),
            key=lambda item: item[0],
        )
    ]
    return ordered_points[:limit]


def compute_inverse_document_frequency(documents: list[list[str]]) -> dict[str, float]:
    total_documents = len(documents)
    document_frequency: Counter[str] = Counter()

    for document in documents:
        document_frequency.update(set(document))

    return {
        term: math.log((1 + total_documents) / (1 + frequency)) + 1.0
        for term, frequency in document_frequency.items()
    }


def build_tfidf_counter(text: str, idf: dict[str, float]) -> Counter[str]:
    term_frequency = build_keyword_counter(text)
    return Counter(
        {
            term: round(count * idf.get(term, 1.0), 4)
            for term, count in term_frequency.items()
        }
    )


def extract_text_from_presentation_bytes(contents: bytes) -> dict[str, object]:
    if not contents:
        raise ValueError("The uploaded presentation was empty.")

    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise RuntimeError(
            "PowerPoint support is not installed yet. Add python-pptx to enable lesson deck uploads."
        ) from exc

    try:
        presentation = Presentation(io.BytesIO(contents))
    except Exception as exc:
        raise ValueError("StudyGenie could not read that PowerPoint file. Upload a .pptx deck.") from exc

    slides: list[dict[str, object]] = []
    all_text_parts: list[str] = []
    documents: list[list[str]] = []

    for index, slide in enumerate(presentation.slides, start=1):
        title_shape = getattr(slide.shapes, "title", None)
        title = clean_fragment(
            title_shape.text if title_shape and getattr(title_shape, "text", "") else ""
        )
        bullet_points: list[str] = []

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            raw_text = normalize_text(getattr(shape, "text", ""))
            if not raw_text or clean_fragment(raw_text) == title:
                continue

            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = clean_fragment(paragraph.text)
                if not paragraph_text or paragraph_text == title:
                    continue
                bullet_points.append(paragraph_text)

        bullet_points = merge_similar_sentences(dedupe_preserve_order(bullet_points), threshold=0.78)
        if not title and not bullet_points:
            continue

        slide_title = title or f"Slide {index}"
        slide_text = normalize_text("\n".join([slide_title, *bullet_points]))
        slide_sentences = merge_similar_sentences(
            split_sentences("\n".join([slide_title, *bullet_points])),
            threshold=0.78,
        )
        slide_keywords = get_ranked_keywords(slide_text, topic=slide_title, limit=8)
        slide_tokens = filtered_tokens(slide_text)

        slides.append(
            {
                "slide_number": index,
                "title": slide_title,
                "points": bullet_points[:10],
                "sentences": slide_sentences,
                "text": slide_text,
                "title_tokens": set(filtered_tokens(slide_title)),
                "keyword_counter": build_keyword_counter(slide_text),
                "keywords": slide_keywords,
                "tfidf_counter": Counter(),
            }
        )
        all_text_parts.append(slide_text)
        documents.append(slide_tokens)

    combined_text = normalize_text("\n".join(all_text_parts))
    if not combined_text or not slides:
        raise ValueError("No readable lesson text was found in that PowerPoint deck.")

    idf = compute_inverse_document_frequency(documents)
    for slide in slides:
        slide["tfidf_counter"] = build_tfidf_counter(str(slide["text"]), idf)

    return {
        "slides": slides,
        "combined_text": combined_text,
        "slide_count": len(slides),
        "idf": idf,
    }


def create_cluster_from_slide(slide: dict[str, object]) -> dict[str, object]:
    return {
        "slides": [slide],
        "title_tokens": set(slide["title_tokens"]),
        "keyword_counter": Counter(slide["keyword_counter"]),
        "tfidf_counter": Counter(slide["tfidf_counter"]),
        "keywords": list(slide["keywords"]),
    }


def append_slide_to_cluster(cluster: dict[str, object], slide: dict[str, object]) -> dict[str, object]:
    cluster["slides"].append(slide)
    cluster["title_tokens"] |= slide["title_tokens"]
    cluster["keyword_counter"].update(slide["keyword_counter"])
    cluster["tfidf_counter"].update(slide["tfidf_counter"])
    cluster["keywords"] = dedupe_preserve_order([*cluster["keywords"], *slide["keywords"]])[:14]
    return cluster


def cluster_reference_slide(cluster: dict[str, object]) -> dict[str, object]:
    return cluster["slides"][-1]


def compute_cluster_similarity(cluster: dict[str, object], slide: dict[str, object]) -> float:
    title_similarity = jaccard_similarity(cluster["title_tokens"], slide["title_tokens"])
    keyword_overlap = jaccard_similarity(set(cluster["keywords"]), set(slide["keywords"]))
    keyword_similarity = cosine_similarity(cluster["keyword_counter"], slide["keyword_counter"])
    tfidf_similarity = cosine_similarity(cluster["tfidf_counter"], slide["tfidf_counter"])
    slide_gap = int(slide["slide_number"]) - int(cluster_reference_slide(cluster)["slide_number"])
    adjacency_bonus = 0.04 if 0 < slide_gap <= 2 else 0.0
    return (
        tfidf_similarity * 0.42
        + keyword_similarity * 0.22
        + title_similarity * 0.2
        + keyword_overlap * 0.16
        + adjacency_bonus
    )


def compute_cluster_pair_similarity(left: dict[str, object], right: dict[str, object]) -> float:
    title_similarity = jaccard_similarity(left["title_tokens"], right["title_tokens"])
    keyword_overlap = jaccard_similarity(set(left["keywords"]), set(right["keywords"]))
    keyword_similarity = cosine_similarity(left["keyword_counter"], right["keyword_counter"])
    tfidf_similarity = cosine_similarity(left["tfidf_counter"], right["tfidf_counter"])
    slide_gap = int(right["slides"][0]["slide_number"]) - int(left["slides"][-1]["slide_number"])
    adjacency_bonus = 0.05 if 0 <= slide_gap <= 2 else 0.0
    return (
        tfidf_similarity * 0.44
        + keyword_similarity * 0.22
        + title_similarity * 0.18
        + keyword_overlap * 0.16
        + adjacency_bonus
    )


def merge_clusters(left: dict[str, object], right: dict[str, object]) -> dict[str, object]:
    merged = {
        "slides": [*left["slides"], *right["slides"]],
        "title_tokens": set(left["title_tokens"]) | set(right["title_tokens"]),
        "keyword_counter": Counter(left["keyword_counter"]),
        "tfidf_counter": Counter(left["tfidf_counter"]),
        "keywords": dedupe_preserve_order([*left["keywords"], *right["keywords"]])[:14],
    }
    merged["keyword_counter"].update(right["keyword_counter"])
    merged["tfidf_counter"].update(right["tfidf_counter"])
    merged["slides"] = sorted(merged["slides"], key=lambda slide: int(slide["slide_number"]))
    return merged


def cluster_slides(slides: list[dict[str, object]]) -> list[dict[str, object]]:
    if not slides:
        return []

    clusters: list[dict[str, object]] = []

    for slide in slides:
        if not clusters:
            clusters.append(create_cluster_from_slide(slide))
            continue

        scored_matches = [
            (
                index,
                compute_cluster_similarity(cluster, slide),
                len(set(cluster["keywords"]) & set(slide["keywords"])),
                len(cluster["title_tokens"] & slide["title_tokens"]),
            )
            for index, cluster in enumerate(clusters)
        ]
        best_index, best_score, shared_keywords, shared_title_terms = max(
            scored_matches,
            key=lambda item: item[1],
        )

        should_join = (
            best_score >= 0.29
            or (best_score >= 0.23 and shared_keywords >= 2)
            or (best_score >= 0.21 and shared_title_terms >= 1)
        )

        if should_join:
            append_slide_to_cluster(clusters[best_index], slide)
        else:
            clusters.append(create_cluster_from_slide(slide))

    clusters = sorted(clusters, key=lambda cluster: int(cluster["slides"][0]["slide_number"]))
    merged_clusters: list[dict[str, object]] = []

    for cluster in clusters:
        if not merged_clusters:
            merged_clusters.append(cluster)
            continue

        previous_cluster = merged_clusters[-1]
        similarity = compute_cluster_pair_similarity(previous_cluster, cluster)
        should_merge = (
            similarity >= 0.32
            or (similarity >= 0.27 and min(len(previous_cluster["slides"]), len(cluster["slides"])) == 1)
        )

        if should_merge:
            merged_clusters[-1] = merge_clusters(previous_cluster, cluster)
        else:
            merged_clusters.append(cluster)

    return merged_clusters


def choose_cluster_title(cluster: dict[str, object]) -> str:
    slide_titles = [clean_fragment(slide["title"]) for slide in cluster["slides"] if slide["title"]]
    non_generic_titles = [title for title in slide_titles if not is_generic_title(title)]
    ranked_titles = non_generic_titles or slide_titles

    if not ranked_titles:
        first_slide = cluster["slides"][0]
        return f"Topic from slide {first_slide['slide_number']}"

    title_counts = Counter(ranked_titles)
    cluster_keywords = set(cluster["keywords"][:6])
    return max(
        ranked_titles,
        key=lambda title: (
            title_counts[title],
            len(set(filtered_tokens(title)) & cluster_keywords),
            len(filtered_tokens(title)),
            len(title),
        ),
    )


def extract_subtopics_from_cluster(cluster: dict[str, object], title: str, limit: int = 4) -> list[str]:
    subtopics: list[str] = []

    for slide in cluster["slides"]:
        candidate = clean_fragment(slide["title"])
        if not candidate or is_generic_title(candidate):
            continue
        if sentence_similarity(candidate, title) >= 0.82:
            continue
        if any(sentence_similarity(candidate, existing) >= 0.82 for existing in subtopics):
            continue
        subtopics.append(candidate)
        if len(subtopics) >= limit:
            break

    if subtopics:
        return subtopics

    fallback_points = merge_similar_sentences(
        [
            point
            for slide in cluster["slides"]
            for point in slide["points"]
        ],
        threshold=0.78,
    )
    fallback_subtopics: list[str] = []
    for point in fallback_points:
        compressed = compress_sentence(point, topic=title, max_words=5)
        if len(filtered_tokens(compressed)) < 2:
            continue
        if any(sentence_similarity(compressed, existing) >= 0.8 for existing in fallback_subtopics):
            continue
        fallback_subtopics.append(compressed.title())
        if len(fallback_subtopics) >= limit:
            break

    return fallback_subtopics


def choose_focus_terms(
    cluster: dict[str, object],
    title: str,
    subtopics: list[str],
    limit: int = 5,
) -> list[str]:
    excluded_terms = set(filtered_tokens(title))
    for subtopic in subtopics:
        excluded_terms.update(filtered_tokens(subtopic))

    focus_terms: list[str] = []
    for keyword in cluster["keywords"]:
        if keyword in excluded_terms or len(keyword) < 4:
            continue
        if keyword in focus_terms:
            continue
        focus_terms.append(keyword)
        if len(focus_terms) >= limit:
            break

    return focus_terms


def estimate_difficulty(summary: str, key_points: list[str], keywords: list[str]) -> str:
    all_words = filtered_tokens(" ".join([summary, *key_points]))
    if not all_words:
        return "easy"

    long_word_ratio = sum(1 for word in all_words if len(word) >= 9) / len(all_words)
    average_words_per_point = sum(len(filtered_tokens(point)) for point in key_points) / max(len(key_points), 1)
    technical_keyword_ratio = sum(1 for word in keywords if len(word) >= 8) / max(len(keywords), 1)
    difficulty_score = long_word_ratio * 2.1 + average_words_per_point / 12 + technical_keyword_ratio

    if difficulty_score >= 1.7:
        return "hard"
    if difficulty_score >= 1.0:
        return "medium"
    return "easy"


def estimate_importance(
    cluster: dict[str, object],
    focus_terms: list[str],
    key_points: list[str],
    deck_keyword_counter: Counter[str],
    total_slides: int,
) -> float:
    repeated_terms = sum(count - 1 for _term, count in cluster["keyword_counter"].items() if count > 1)
    cluster_term_total = max(1, sum(cluster["keyword_counter"].values()))
    deck_total = max(1, sum(deck_keyword_counter.values()))
    deck_focus_mass = sum(deck_keyword_counter.get(term, 0) for term in set(focus_terms))
    coverage_score = min(1.0, len(cluster["slides"]) / max(1, total_slides * 0.35))
    repetition_score = min(1.0, repeated_terms / max(1, cluster_term_total * 0.22))
    deck_score = min(1.0, deck_focus_mass / max(1, deck_total * 0.08))
    key_point_score = min(1.0, len(key_points) / 4)

    score = (
        0.18
        + coverage_score * 0.2
        + repetition_score * 0.3
        + deck_score * 0.24
        + key_point_score * 0.08
        + min(1.0, len(focus_terms) / 5) * 0.06
    )
    return round(max(0.05, min(0.99, score)), 2)


def build_section_summary(
    title: str,
    sentences: list[str],
    points: list[str],
    focus_terms: list[str],
    tfidf_weights: Counter[str],
) -> str:
    candidate_pool = merge_similar_sentences([*points, *sentences], threshold=0.72)
    ranked_summary = build_ranked_summary(
        candidate_pool,
        topic=title,
        max_sentences=2,
        tfidf_weights=tfidf_weights,
    )
    selected = split_sentences(ranked_summary) if ranked_summary else []

    compressed_sentences: list[str] = []
    for sentence in selected or candidate_pool[:2]:
        compressed = compress_sentence(sentence, topic=title, max_words=14)
        if not compressed:
            continue
        if any(sentence_similarity(compressed, existing) >= 0.82 for existing in compressed_sentences):
            continue
        compressed_sentences.append(ensure_sentence(compressed))
        if len(compressed_sentences) >= 2:
            break

    if compressed_sentences:
        return shorten_phrase(" ".join(compressed_sentences), max_words=30)

    if focus_terms:
        return ensure_sentence(shorten_phrase(f"{title} focuses on {', '.join(focus_terms[:3])}", max_words=18))

    return ensure_sentence(title)


def build_section_from_cluster(
    cluster: dict[str, object],
    deck_keyword_counter: Counter[str],
    total_slides: int,
) -> dict[str, object]:
    title = choose_cluster_title(cluster)

    sentences: list[str] = []
    points: list[str] = []
    slide_numbers: list[int] = []
    tfidf_weights: Counter[str] = Counter()

    for slide in cluster["slides"]:
        sentences.extend(slide["sentences"])
        points.extend(slide["points"])
        slide_numbers.append(int(slide["slide_number"]))
        tfidf_weights.update(slide["tfidf_counter"])

    merged_sentences = merge_similar_sentences(sentences, threshold=0.74)
    merged_points = merge_similar_sentences(points, threshold=0.78)
    subtopics = extract_subtopics_from_cluster(cluster, title)
    focus_terms = choose_focus_terms(cluster, title, subtopics)
    key_points = build_key_points(
        [*merged_points, *merged_sentences],
        topic=title,
        limit=4,
        tfidf_weights=tfidf_weights,
    )
    summary = build_section_summary(title, merged_sentences, merged_points, focus_terms, tfidf_weights)
    difficulty = estimate_difficulty(summary, key_points, focus_terms)
    importance = estimate_importance(
        cluster,
        focus_terms=focus_terms,
        key_points=key_points,
        deck_keyword_counter=deck_keyword_counter,
        total_slides=total_slides,
    )

    return {
        "title": title,
        "subtopics": subtopics,
        "summary": summary,
        "key_points": key_points,
        "difficulty": difficulty,
        "importance": importance,
        "focus_terms": focus_terms,
        "slide_numbers": slide_numbers,
    }


def build_overview(
    lesson_title: str,
    sections: list[dict[str, object]],
    all_sentences: list[str],
    keywords: list[str],
) -> str:
    overview_candidates = [section["summary"] for section in sections[:4]] + all_sentences[:8]
    ranked = build_ranked_summary(overview_candidates, topic=lesson_title, max_sentences=2)
    compressed = [
        ensure_sentence(compress_sentence(sentence, topic=lesson_title, max_words=16))
        for sentence in split_sentences(ranked)
    ]
    compressed = [sentence for sentence in compressed if sentence]

    if compressed:
        return shorten_phrase(" ".join(compressed[:2]), max_words=34)

    section_titles = ", ".join(section["title"] for section in sections[:3])
    if section_titles:
        fallback = f"{lesson_title} is organized around {section_titles}"
        if keywords:
            fallback += f" and repeated ideas such as {', '.join(keywords[:3])}"
        return ensure_sentence(shorten_phrase(fallback, max_words=32))

    return ensure_sentence(shorten_phrase(lesson_title, max_words=10))


def estimate_revision_time(sections: list[dict[str, object]]) -> str:
    total_minutes = 6
    for section in sections:
        total_minutes += 2
        total_minutes += len(section["key_points"]) * 2
        if section["difficulty"] == "hard":
            total_minutes += 4
        elif section["difficulty"] == "medium":
            total_minutes += 2
        if float(section["importance"]) >= 0.75:
            total_minutes += 2

    rounded = max(10, int(math.ceil(total_minutes / 5.0) * 5))
    return f"{rounded} mins"


def build_quiz_questions(sections: list[dict[str, object]]) -> list[str]:
    questions: list[str] = []

    for section in sections[:4]:
        if section["subtopics"]:
            questions.append(
                f"How do {section['subtopics'][0]} and {section['title']} connect in the lesson?"
            )
        elif section["key_points"]:
            first_point = section["key_points"][0].rstrip(".")
            questions.append(
                f"How would you explain {section['title']} using the idea: {first_point}?"
            )
        else:
            questions.append(
                f"What is the main idea behind {section['title']} and where would you apply it?"
            )

    return dedupe_preserve_order(questions)[:4]


def build_flashcards(sections: list[dict[str, object]]) -> list[dict[str, str]]:
    flashcards: list[dict[str, str]] = []

    for section in sections[:4]:
        flashcards.append(
            {
                "front": f"What should you remember about {section['title']}?",
                "back": section["summary"],
            }
        )

        if section["key_points"]:
            flashcards.append(
                {
                    "front": f"Key point from {section['title']}",
                    "back": section["key_points"][0],
                }
            )

    return flashcards[:6]


def slugify_concept_key(text: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", clean_fragment(text).lower()).strip("-")
    return slug or "concept"


def build_concept_graph(
    sections: list[dict[str, object]],
) -> tuple[list[dict[str, object]], list[dict[str, object]]]:
    concepts: list[dict[str, object]] = []
    concept_edges: list[dict[str, object]] = []
    seen_keys: set[str] = set()
    previous_section: dict[str, object] | None = None

    for section_index, section in enumerate(sections, start=1):
        section_title = clean_fragment(str(section["title"]))
        section_key = f"section-{section_index}-{slugify_concept_key(section_title)}"
        while section_key in seen_keys:
            section_key = f"{section_key}-alt"
        seen_keys.add(section_key)

        subtopics = dedupe_preserve_order([str(item) for item in section.get("subtopics", [])])
        focus_terms = dedupe_preserve_order([str(item) for item in section.get("focus_terms", [])])[:5]
        slide_numbers = [int(value) for value in section.get("slide_numbers", [])]
        section_node = {
            "concept_key": section_key,
            "name": section_title,
            "kind": "section",
            "parent_name": None,
            "summary": ensure_sentence(str(section["summary"])),
            "difficulty": str(section["difficulty"]),
            "importance": float(section["importance"]),
            "focus_terms": focus_terms,
            "slide_numbers": slide_numbers,
            "related_concepts": subtopics[:4],
        }
        concepts.append(section_node)

        if previous_section is not None:
            concept_edges.append(
                {
                    "source_concept_key": str(previous_section["concept_key"]),
                    "target_concept_key": section_key,
                    "source_name": str(previous_section["name"]),
                    "target_name": section_title,
                    "relation_type": "progression",
                    "weight": round(
                        0.34
                        + (
                            float(previous_section["importance"])
                            + float(section_node["importance"])
                        )
                        / 4,
                        2,
                    ),
                }
            )
            previous_related = dedupe_preserve_order(
                [*previous_section["related_concepts"], section_title]
            )[:5]
            previous_section["related_concepts"] = previous_related
            section_node["related_concepts"] = dedupe_preserve_order(
                [*section_node["related_concepts"], str(previous_section["name"])]
            )[:5]

        previous_section = section_node
        previous_subtopic: dict[str, object] | None = None

        for subtopic_index, subtopic in enumerate(subtopics, start=1):
            subtopic_key = f"{section_key}-sub-{subtopic_index}-{slugify_concept_key(subtopic)}"
            while subtopic_key in seen_keys:
                subtopic_key = f"{subtopic_key}-alt"
            seen_keys.add(subtopic_key)

            related_concepts = [section_title]
            if previous_subtopic is not None:
                related_concepts.append(str(previous_subtopic["name"]))

            concept_summary = compress_sentence(str(section["summary"]), topic=subtopic, max_words=18)
            if not concept_summary and section.get("key_points"):
                concept_summary = compress_sentence(
                    str(section["key_points"][0]),
                    topic=subtopic,
                    max_words=18,
                )

            subtopic_node = {
                "concept_key": subtopic_key,
                "name": subtopic,
                "kind": "subtopic",
                "parent_name": section_title,
                "summary": ensure_sentence(
                    concept_summary or f"{subtopic} supports the wider idea in {section_title}"
                ),
                "difficulty": str(section["difficulty"]),
                "importance": round(max(0.08, float(section["importance"]) * 0.88), 2),
                "focus_terms": focus_terms,
                "slide_numbers": slide_numbers,
                "related_concepts": dedupe_preserve_order(related_concepts)[:4],
            }
            concepts.append(subtopic_node)

            concept_edges.append(
                {
                    "source_concept_key": section_key,
                    "target_concept_key": subtopic_key,
                    "source_name": section_title,
                    "target_name": subtopic,
                    "relation_type": "contains",
                    "weight": round(max(0.54, float(section["importance"])), 2),
                }
            )

            if previous_subtopic is not None:
                concept_edges.append(
                    {
                        "source_concept_key": str(previous_subtopic["concept_key"]),
                        "target_concept_key": subtopic_key,
                        "source_name": str(previous_subtopic["name"]),
                        "target_name": subtopic,
                        "relation_type": "supports",
                        "weight": 0.48,
                    }
                )

            previous_subtopic = subtopic_node

    return concepts, concept_edges


def build_lesson_summary_from_presentation(contents: bytes) -> dict[str, object]:
    extracted = extract_text_from_presentation_bytes(contents)
    slides = extracted["slides"]
    combined_text = extracted["combined_text"]
    slide_count = extracted["slide_count"]

    lesson_title = next(
        (str(slide["title"]) for slide in slides if not is_generic_title(str(slide["title"]))),
        str(slides[0]["title"]),
    )

    all_sentences = merge_similar_sentences(split_sentences(combined_text), threshold=0.76)
    deck_keyword_counter = build_keyword_counter(combined_text, excluded=set(filtered_tokens(lesson_title)))
    clusters = cluster_slides(slides)
    sections = [
        build_section_from_cluster(
            cluster,
            deck_keyword_counter=deck_keyword_counter,
            total_slides=slide_count,
        )
        for cluster in clusters
    ]
    sections = sorted(
        sections,
        key=lambda section: min(section["slide_numbers"]) if section["slide_numbers"] else 10_000,
    )

    keywords = get_ranked_keywords(combined_text, topic=lesson_title, limit=10)
    overview = build_overview(lesson_title, sections, all_sentences, keywords)
    revise_first = sorted(
        sections,
        key=lambda section: (section["importance"], len(section["key_points"])),
        reverse=True,
    )[:3]
    estimated_revision_time = estimate_revision_time(sections)
    quiz_questions = build_quiz_questions(revise_first or sections)
    flashcards = build_flashcards(revise_first or sections)
    concepts, concept_edges = build_concept_graph(sections)

    return {
        "title": lesson_title,
        "overview": overview,
        "keywords": keywords,
        "slides": [
            {
                "slide_number": int(slide["slide_number"]),
                "title": str(slide["title"]),
                "points": [str(point) for point in slide.get("points", [])],
                "text": str(slide["text"]),
            }
            for slide in slides
        ],
        "sections": [
            {
                "title": section["title"],
                "subtopics": section["subtopics"],
                "summary": section["summary"],
                "key_points": section["key_points"],
                "difficulty": section["difficulty"],
                "importance": section["importance"],
                "focus_terms": section["focus_terms"],
                "slide_numbers": section["slide_numbers"],
            }
            for section in sections
        ],
        "revise_first": [
            {
                "title": section["title"],
                "subtopics": section["subtopics"],
                "summary": section["summary"],
                "key_points": section["key_points"],
                "difficulty": section["difficulty"],
                "importance": section["importance"],
                "focus_terms": section["focus_terms"],
                "slide_numbers": section["slide_numbers"],
            }
            for section in revise_first
        ],
        "quiz_questions": quiz_questions,
        "flashcards": flashcards,
        "concepts": concepts,
        "concept_edges": concept_edges,
        "estimated_revision_time": estimated_revision_time,
        "slide_count": slide_count,
        "source_text_length": len(combined_text),
    }
